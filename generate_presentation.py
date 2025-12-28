from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # Helper to add a title slide
    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[0] # 0 is usually Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle

    # Helper to add a content slide (Title + Bullets)
    def add_content_slide(title_text, content_items):
        slide_layout = prs.slide_layouts[1] # 1 is Title + Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.font.size = Pt(20)

    # Helper to add code block slide (Special layout or just text)
    def add_code_slide(title_text, code_text):
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        p.text = code_text
        p.font.name = 'Courier New'
        p.font.size = Pt(12)

    # --- SLIDES ---

    # Slide 1
    add_title_slide("EMOVISTA", "Real-Time Multimodal Emotion-Aware Assistant")

    # Slide 2
    add_content_slide("Problem Statement", [
        "Understanding human emotions in real-time is critical for therapeutic, safety, and interactive AI applications.",
        "Existing solutions often rely on a single modality (vision or audio) and require cloud services."
    ])

    # Slide 3
    add_content_slide("Key Features", [
        "Multimodal Fusion: Combines facial expression (FER), speech emotion (SER), and text sentiment.",
        "Offline-First: All models run locally, no external API calls.",
        "Medical-Ready: Encrypted patient memory, severity scoring, emergency escalation.",
        "Rich UI: Desktop (CustomTkinter) and web (Flask) interfaces with dark-mode aesthetics."
    ])

    # Slide 4
    add_code_slide("Architecture Overview", 
"""+-------------------+      +-------------------+
|   Front-End UI    | ---> |   Flask Backend   |
+-------------------+      +-------------------+
             |                       |
             v                       v
          +-------------------------------+
          |   Fusion Engine (emotion_fusion) |
          +-------------------------------+
            /        |          \        \\
           v         v           v        v
     FER Model   Speech Model   Text Model   Vosk STT
     
- All models are stored under models/
- Vosk provides offline speech-to-text
- tts.py gives spoken feedback""")

    # Slide 5
    add_content_slide("Vosk Model Update", [
        "Switched to vosk-model-en-us-0.22 for higher accuracy.",
        "Updated app.py and README accordingly."
    ])

    # Slide 6
    add_content_slide("Text Emotion Fusion", [
        "/consult endpoint now:",
        "1. Predicts text sentiment using text_model.",
        "2. Converts frontend FER label to a one-hot probability vector.",
        "3. Calls fuse() to obtain a fused emotion.",
        "4. Generates assistant response based on the fused label."
    ])

    # Slide 7
    add_content_slide("Startup Greeting", [
        "On server start, the system announces: 'Emovista system online and ready.'",
        "Uses tts.speak() with a guard (WERKZEUG_RUN_MAIN) to avoid double-speaking during Flask reload."
    ])

    # Slide 8
    add_content_slide("Demo Flow (Web UI)", [
        "1. User opens http://localhost:5000",
        "2. Webcam captures face -> FER prediction.",
        "3. User types or speaks text -> text sentiment.",
        "4. Fusion engine produces final emotion.",
        "5. Assistant replies with voice feedback."
    ])

    # Slide 9
    add_content_slide("Security & Privacy", [
        "Patient data encrypted with AES (patient_memory.py).",
        "No cloud calls – all processing stays on-device.",
        "Severity engine flags high-risk emotional states."
    ])

    # Slide 10
    add_content_slide("Practical Use in the Medical Field", [
        "Patient Monitoring: Real-time emotional cues during tele-health sessions.",
        "Therapeutic Feedback: Detect distress/anxiety and trigger calming interventions.",
        "Emergency Escalation: Alerts caregivers when negative emotions persist.",
        "Privacy-First: Local processing and encryption meet HIPAA-like requirements.",
        "Integration: Backend can be called from EHR systems via REST APIs."
    ])

    # Slide 11
    add_content_slide("Future Work", [
        "Train a learned fusion model for higher accuracy.",
        "Add multilingual Vosk models.",
        "Integrate more expressive TTS voices.",
        "Deploy as a Docker container for easy distribution."
    ])

    # Slide 12
    add_code_slide("Get Started", 
"""# Clone repo
git clone https://github.com/nohanbaloch/EMOVISTA.git
cd EMOVISTA

# Install dependencies
pip install -r requirements.txt

# Download Vosk model
wget https://alphacephei.com/vosk/models/vosk-model-en-us-0.22.zip
unzip vosk-model-en-us-0.22.zip -d models/vosk/

# Run the web backend
python src/web/backend/app.py""")

    # Slide 13
    add_title_slide("Thank You", "Questions? Contact: nohan@example.com")

    # Save
    output_path = "EMOVISTA_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
